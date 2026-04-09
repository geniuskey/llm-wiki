import os
import subprocess
from fastapi import FastAPI, File, UploadFile, Form, BackgroundTasks
from fastapi.middleware.cors import CORSMiddleware
from dotenv import load_dotenv
import shutil
from datetime import datetime
import pdfplumber
from docx import Document
from pptx import Presentation
import openpyxl
import requests
from bs4 import BeautifulSoup
from typing import Optional, Dict, Any
from pydantic import BaseModel

from my_rag import RagClient

load_dotenv(dotenv_path="../.env")

app = FastAPI()
rag = RagClient()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

RAW_SOURCES_DIR = "../wiki/raw_sources"
os.makedirs(RAW_SOURCES_DIR, exist_ok=True)

def parse_file_to_markdown(filepath: str, filename: str) -> str:
    text_content = ""
    try:
        ext = filename.lower().split('.')[-1]
        if ext == 'pdf':
            with pdfplumber.open(filepath) as pdf:
                text_content = "\n".join([page.extract_text() or '' for page in pdf.pages])
        elif ext == 'docx':
            doc = Document(filepath)
            text_content = "\n".join([p.text for p in doc.paragraphs])
        elif ext == 'pptx':
            prs = Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"
        elif ext == 'xlsx':
            wb = openpyxl.load_workbook(filepath, data_only=True)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_data = [str(cell) for cell in row if cell is not None]
                    if row_data:
                        text_content += " | ".join(row_data) + "\n"
        elif ext in ['md', 'txt']:
            with open(filepath, 'r', encoding='utf-8') as f:
                text_content = f.read()
        else:
            text_content = f"Unsupported file type: {ext}"
    except Exception as e:
        text_content = f"Error parsing file: {e}"
        
    return f"# {filename}\n\n{text_content}"

def trigger_claude_ingest(filename: str):
    print(f"Triggering Claude for {filename}...")
    prompt = f"A new source has been added to ../raw_sources/{filename}. Please ingest it, update relevant entity/concept pages, update the docs/index.md, and add an entry to docs/log.md according to our wiki schema. Make sure you fully read the new source."
    
    try:
        # Note: --yes or similar auto-confirm flags might be needed depending on Claude CLI version
        result = subprocess.run(
            ["claude", "-p", prompt],
            cwd="../wiki",
            capture_output=True,
            text=True
        )
        print("Claude output:", result.stdout)
        if result.stderr:
            print("Claude error:", result.stderr)

        # After Claude finishes, check deployment mode
        deployment_mode = os.getenv("DEPLOYMENT_MODE", "local")
        if deployment_mode == "github":
            print("Deploying via GitHub...")
            # We are running inside backend/, so we reference wiki/ with ../wiki
            subprocess.run(["git", "add", "."], cwd="../wiki")
            subprocess.run(["git", "commit", "-m", f"Auto-update wiki for {filename}"], cwd="../wiki")
            subprocess.run(["git", "push", "origin", "main"], cwd="../wiki")
        else:
            print("Deploying locally via MkDocs build...")
            subprocess.run(["mkdocs", "build"], cwd="../wiki")
            
    except Exception as e:
        print(f"Error executing claude or deployment: {e}")

@app.post("/upload")
async def upload_file(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    # Save the file temporarily
    temp_dir = "../temp_uploads"
    os.makedirs(temp_dir, exist_ok=True)
    temp_path = os.path.join(temp_dir, file.filename)
    
    with open(temp_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    # Convert file to markdown
    md_filename = f"{datetime.now().strftime('%Y%m%d%H%M%S')}_{file.filename}.md"
    md_filepath = os.path.join(RAW_SOURCES_DIR, md_filename)
    
    markdown_content = parse_file_to_markdown(temp_path, file.filename)
    with open(md_filepath, "w", encoding="utf-8") as f:
        f.write(markdown_content)
        
    # RAG 시스템에 문서 인덱싱
    rag.insert_doc(
        doc_id=md_filename,
        content=markdown_content,
        metadata={"source_type": "upload", "original_filename": file.filename}
    )
        
    # Queue parsing background task
    background_tasks.add_task(trigger_claude_ingest, md_filename)
    
    return {"message": "File processed and queued for wiki ingestion.", "filename": md_filename}

@app.post("/url")
async def submit_url(background_tasks: BackgroundTasks, url: str = Form(...)):
    try:
        # Basic Generic Fetching for testing
        # In a real environment, you'd use Jira/Confluence API with the tokens from .env
        headers = {'User-Agent': 'Mozilla/5.0'}
        res = requests.get(url, headers=headers, timeout=10)
        res.raise_for_status()
        soup = BeautifulSoup(res.text, 'html.parser')
        text_content = soup.get_text(separator='\n', strip=True)
    except Exception as e:
        text_content = f"Failed to fetch URL: {e}"
        
    md_filename = f"url_{datetime.now().strftime('%Y%m%d%H%M%S')}.md"
    md_filepath = os.path.join(RAW_SOURCES_DIR, md_filename)
    
    markdown_content = f"# Extracted from URL: {url}\n\n{text_content}"
    with open(md_filepath, "w", encoding="utf-8") as f:
        f.write(markdown_content)
        
    # RAG 시스템에 문서 인덱싱
    rag.insert_doc(
        doc_id=md_filename,
        content=markdown_content,
        metadata={"source_type": "url", "url": url}
    )
        
    background_tasks.add_task(trigger_claude_ingest, md_filename)
    
    return {"message": "URL processed and queued for wiki ingestion.", "filename": md_filename}

class AskRequest(BaseModel):
    query: str
    filters: Optional[Dict[str, Any]] = None

@app.post("/ask")
async def ask_rag(request: AskRequest):
    """
    RAG 시스템을 사용하여 질문에 답변을 제공하기 위한 엔드포인트
    (현재는 검색된 문서 결과만 반환)
    """
    results = rag.retrieve(query=request.query, filters=request.filters)
    
    return {
        "query": request.query,
        "retrieved_documents": results,
        "message": "추후 LLM API를 연동하여 검색된 문서를 바탕으로 최종 답변을 생성하도록 확장할 수 있습니다."
    }
